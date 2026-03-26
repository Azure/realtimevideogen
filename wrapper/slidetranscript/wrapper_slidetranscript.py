"""
Slide transcript generation wrapper.
"""

import logging
import json
import tempfile
import aiofiles

from openai import AsyncOpenAI

from typing import TYPE_CHECKING
from typing import cast
from typing import override
from typing import Dict
from typing import Optional
from typing import List
from typing import Tuple
from typing import Any
from typing import AsyncGenerator

if TYPE_CHECKING:
    from openai import AsyncStream
    from openai.types.chat import ChatCompletionChunk

from wrapper_model import ModelGeneration

from file_utils import base64_to_binary
from media_utils import fix_json_like_string
from file_utils import binary_to_base64

from transcript_prompts import SYSTEM_PROMPT

from pptx import Presentation

from ppt_utils import pptx_to_images


DEFAULT_MAX_TOKENS = 5000
DEFAULT_TEMPERATURE = 0.7


class SlideTranscriptGenerator(ModelGeneration):
    """A wrapper for slide transcript generation using an LLM."""

    def __init__(
        self,
        llm_url: str = "http://localhost:8000/v1",
        llm_model: str = "meta-llama/Meta-Llama-3.1-8B",
    ) -> None:
        super().__init__("slidetranscript")
        self.set_llm(llm_model, llm_url)

    def set_llm(
        self,
        llm_model: str,
        llm_url: str,
        api_key: str = "n/a",
    ) -> None:
        self.llm_model = llm_model
        self.llm_url = llm_url
        self.llm_client = AsyncOpenAI(
            api_key=api_key,
            base_url=self.llm_url)
        self.extra_body = dict(guided_decoding_backend="xgrammar")

    def parse_pptx(
        self,
        pptx_path: str
    ) -> Tuple[Optional[List[str]], Optional[List[str]]]:
        pptx_texts = []
        pptx_images: List[str] = []

        # Add the text from each slide
        presentation = Presentation(pptx_path)
        slide_ix = 0
        for slide in presentation.slides:
            is_hidden = slide._element.get("show") == "0"
            if not is_hidden:
                slide_text = [
                    f"--- SLIDE {slide_ix + 1} ---"
                ]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            slide_text.append(paragraph.text)
                pptx_texts.append("\n".join(slide_text))
                slide_ix += 1
        logging.info(f"Detected {slide_ix} slides.")
        logging.info(f"Extracted text from {len(pptx_texts)} slides.")

        # Render each of the slides
        pptx_image_paths = pptx_to_images(
            pptx_path,
            output_path="/tmp",
        )
        for pptx_image_path in pptx_image_paths:
            with open(pptx_image_path, "rb") as image_file:
                image_binary = image_file.read()
                image_base64 = binary_to_base64(image_binary)
                pptx_images.append(image_base64)
        logging.info(f"Rendered {len(pptx_images)} slide images.")

        return pptx_texts, pptx_images

    @override
    async def generate(
        self,
        pptx_texts: Optional[List[str]] = None,
        pptx_images: Optional[List[str]] = None,
        max_tokens: int = DEFAULT_MAX_TOKENS,
        temperature: float = DEFAULT_TEMPERATURE,
        llm_model: Optional[str] = None,
        llm_url: Optional[str] = None,
        max_words_per_slide: int = -1,
        job_id: Optional[str] = None,
    ) -> List[Dict[str, Any]]:
        gen_timer = self._new_gen_timer(job_id)

        self.running = True  # We can run in parallel but good to know if we are running

        if llm_model is not None and llm_url is not None:
            logging.info(f"Setting LLM model to {llm_model} at {llm_url}.")
            self.set_llm(llm_model, llm_url)

        try:
            ret = []
            async for slide_text in self.generate_stream(
                pptx_texts,
                pptx_images,
                max_tokens=max_tokens,
                temperature=temperature,
                max_words_per_slide=max_words_per_slide,
                job_id=job_id,
            ):
                ret.append(slide_text)
            return ret
        finally:
            gen_timer.end("total")

    async def generate_stream(
        self,
        pptx_texts: Optional[List[str]] = None,
        pptx_images: Optional[List[str]] = None,
        max_tokens: int = DEFAULT_MAX_TOKENS,
        temperature: float = DEFAULT_TEMPERATURE,
        llm_model: Optional[str] = None,
        llm_url: Optional[str] = None,
        max_words_per_slide: int = -1,
        job_id: Optional[str] = None,
    ) -> AsyncGenerator[Dict, None]:
        gen_timer = self._new_gen_timer(job_id)

        self.running = True  # We can run in parallel but good to know if we are running

        if llm_model is not None and llm_url is not None:
            logging.info(f"Setting LLM model to {llm_model} at {llm_url}.")
            self.set_llm(llm_model, llm_url)

        try:
            it = 0
            gen_timer.start("gen_script_stream")
            gen_timer.start(f"gen_script_stream_{it}")
            async for slide_text in self.gen_script_stream(
                pptx_texts,
                pptx_images,
                max_words_per_slide=max_words_per_slide,
                job_id=job_id,
            ):
                yield slide_text
                gen_timer.end(f"gen_script_stream_{it}")
                it += 1
                gen_timer.start(f"gen_script_stream_{it}")
            gen_timer.end(f"gen_script_stream_{it}")
            gen_timer.end("gen_script_stream")
        finally:
            self.running = False
            gen_timer.end("total")

    async def gen_script_stream(
        self,
        pptx_texts: Optional[List[str]] = None,
        pptx_images: Optional[List[str]] = None,
        max_tokens: int = DEFAULT_MAX_TOKENS,
        temperature: float = DEFAULT_TEMPERATURE,
        max_words_per_slide: int = -1,
        job_id: Optional[str] = None,
    ) -> AsyncGenerator[Dict, None]:
        if not pptx_texts:
            raise ValueError("pptx_texts is required to generate slide transcript.")
        if not pptx_images:
            raise ValueError("pptx_images is required to generate slide transcript.")
        if len(pptx_texts) != len(pptx_images):
            raise ValueError(
                f"pptx_texts ({len(pptx_texts)}) and "
                f"pptx_images ({len(pptx_images)}) must have the same length.")

        contraint_instruction = ""
        if max_words_per_slide > 0:
            contraint_instruction += f"Each slide transcript should not exceed {max_words_per_slide} words.\n"

        # Compose prompt
        system_prompt_text = (
            f"{SYSTEM_PROMPT}\n\n"
            f"{contraint_instruction}"
        )

        # Build the messages
        user_msg: List[Dict[str, Any]] = []
        for slide_num, (slide_text, slide_image) in enumerate(zip(pptx_texts, pptx_images)):
            user_msg.append({
                "type": "text",
                "text": (
                    f"{slide_text}\n"
                    "The following image corresponds exactly to this slide.\n"
                )
            })
            user_msg.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/jpeg;base64,{slide_image}"
                }
            })
        user_msg.append({
            "type": "text",
            "text": "Generate the JSONL transcript."
        })

        messages: List[Dict[str, Any]] = [
            {"role": "system", "content": system_prompt_text},
            {"role": "user", "content": user_msg}
        ]

        async with aiofiles.open(f"/tmp/{job_id}_prompt.json", "w") as prompt_file:
            await prompt_file.write(json.dumps(messages, indent=2))

        response_stream = cast("AsyncStream[ChatCompletionChunk]", await self.llm_client.chat.completions.create(
            model=self.llm_model,
            messages=messages,  # type: ignore[arg-type]
            temperature=temperature,
            max_tokens=max_tokens,
            extra_body=self.extra_body,
            stream=True,
        ))

        it = 0
        buffer_text = ""
        async for chunk in response_stream:
            if self.interrupted:  # type: ignore[has-type]
                self.interrupted = False
                logging.info("Generation interrupted.")
                return

            delta = chunk.choices[0].delta.content or ""
            buffer_text += delta
            if delta.endswith("\n"):
                buffer_text = buffer_text.strip()
                if buffer_text.startswith("{") and buffer_text.endswith("}"):
                    try:
                        buffer_text = fix_json_like_string(buffer_text)
                        buffer_json = json.loads(buffer_text)
                        yield buffer_json
                    except json.JSONDecodeError as json_error:
                        logging.error(f"JSON error: {json_error} for buffer: {buffer_text}")
                else:
                    logging.info(f"Ignoring: {buffer_text}")
                buffer_text = ""
            it += 1

    async def warmup(self) -> None:
        pass  # No specific warmup needed for this model

    async def get_rest_args(
        self,
        data_json: Dict[str, str]
    ) -> Dict[str, Any]:
        if data_json is None:
            raise ValueError("Missing JSON body")

        job_id = data_json.get("job_id", None)

        pptx_base64 = data_json.get("pptx", None)
        pptx_texts = None
        pptx_images = None
        if pptx_base64 is not None:
            if not job_id:
                output_path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
            else:
                output_path = f"/tmp/{job_id}.pptx"
            pptx_binary = base64_to_binary(pptx_base64)
            async with aiofiles.open(output_path, "wb") as file:
                await file.write(pptx_binary)
            pptx_texts, pptx_images = self.parse_pptx(output_path)

        rest_args: Dict[str, Any] = {
            "job_id": job_id,
            "llm_url": data_json.get("llm_url", "http://localhost:8000/v1"),
            "llm_model": data_json.get("llm_model", "google/gemma-3-27b-it"),
            "max_words_per_slide": data_json.get("max_words_per_slide", -1)
        }

        if pptx_texts is not None:
            rest_args["pptx_texts"] = pptx_texts
        if pptx_images is not None:
            rest_args["pptx_images"] = pptx_images

        return {
            "task": self.model_name,
            "args": rest_args
        }
