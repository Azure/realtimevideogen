"""
StreamPersona job to generate a video podcast.
It coordinates the execution of the different models.
"""

import sys
import time
import json
import aiofiles
import aiofiles.os
import asyncio
import math

from typing import override
from typing import Dict
from typing import Any
from typing import List
from typing import Tuple
from typing import Optional

from PIL import Image

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from ppt_utils import pptx_to_images
from ppt_utils import get_num_slides

from persona_prompts import IMG_PROMPT
from persona_prompts import IMG_NEG_PROMPT
from persona_prompts import VIDEO_PROMPT
from persona_prompts import VIDEO_NEG_PROMPT


# Local relative imports
sys.path.append("..")  # noqa: E402

from streamwise_job import StreamWiseJob
from streamwise_job import JobStatus
from streamwise_job import OutputMode

from lmm_service_manager import LMMServiceManager

from client import ServiceError

from gen_video_chunked import GenVideoChunked

from tts_utils import estimate_num_words_from_audio_duration

from console_utils import bytes_to_human

from file_utils import read_file_bytes
from file_utils import save_base64_as_binary

from media_utils import get_audio_duration
from media_utils import save_video_audio
from media_utils import get_video_frames
from media_utils import get_video_file_info
from media_utils import get_audio_file_info
from media_utils import get_frame_with_text
from media_utils import concatenate_videos

from k8s_utils import NoActiveContainerError
from k8s_utils import NoRunnableContainerError
from k8s_utils import ServiceNotFoundError

from video import MAX_FT_DURATION_SECS
from video import FANTASYTALKING_FPS


MAX_LOG_TEXT = 100
MAX_IMG_LINE_CHARS = 50


def overlay_image_on_image(
    base_image: Image.Image,
    overlay_image: Image.Image,
    position: Tuple[str, str] = ("bottom", "right"),
    overlay_percentage: float = 0.25
) -> Image.Image:
    """
    Overlay an image on top of another image at the specified position.
    Positions can be 'top', 'bottom', 'left', 'right', 'center'.
    """
    base_width, base_height = base_image.size
    overlay_width, overlay_height = overlay_image.size
    overlay_image = overlay_image.copy().resize((
        int(base_width * overlay_percentage),
        int(base_height * overlay_percentage)
    ))
    overlay_width, overlay_height = overlay_image.size

    if position[0] == "top":
        y = 0
    elif position[0] == "bottom":
        y = base_height - overlay_height
    else:  # center
        y = (base_height - overlay_height) // 2

    if position[1] == "left":
        x = 0
    elif position[1] == "right":
        x = base_width - overlay_width
    else:  # center
        x = (base_width - overlay_width) // 2

    new_image = base_image.copy()
    new_image.paste(
        overlay_image,
        (x, y),
        overlay_image.convert("RGBA"))
    return new_image


class StreamPersonaJob(StreamWiseJob):
    """A job to generate a podcast with images, audio, and video."""

    def __init__(
        self,
        job_id: str,
        service_manager: LMMServiceManager,
        config: Dict[str, Any] = {},
    ) -> None:
        super().__init__(
            "streampersona",
            job_id,
            service_manager,
            config)
        self.image: Optional[Image.Image] = None
        self.image_task: Optional[asyncio.Task[Image.Image]] = None

    @override
    async def generate(
        self,
        job_config: Dict[str, Any],
    ) -> None:
        pptx_base64 = job_config.get("pptx_base64", None)
        if pptx_base64 is None:
            raise ValueError("Missing 'pptx_base64' in job config")
        await self.gen_persona(pptx_base64)

    async def gen_slide_video(
        self,
        slide_number: int,
        slide_text: str,
    ):
        """
        Generate the video for one slide.
        """
        t0 = time.time()
        voice = "am_adam"  # TODO select voice
        audio_base64 = await self.gen.gen_audio(
            slide_text,
            voice=voice,
            task_id=f"{slide_number:03d}",
            deadline=self.get_slide_deadline(slide_number),
        )
        if not audio_base64:
            raise ValueError(f"Cannot generate audio for slide {slide_number} with text '{slide_text}'")
        audio_duration = get_audio_duration(audio_base64)
        audio_path = f"{self.job_path}/{slide_number:03d}.wav"
        audio_len = bytes_to_human(len(audio_base64))
        await save_base64_as_binary(audio_path, audio_base64)
        self.logger.info(
            f"[{slide_number}] Generated audio with {audio_len} and {audio_duration:.3f} seconds.")

        # We need to wait for the image generated
        if self.image is None and self.image_task is not None:
            self.image = await self.image_task
            if self.image is None:
                raise ValueError("Image generation task completed with no image.")
            width, height = self.image.size
            image_path = f"{self.job_path}/persona.png"
            self.image.save(image_path)
            self.logger.info(f"Image with {width}x{height} pixels saved to '{image_path}'.")

        output_mode = self.get_config_output_mode()
        if output_mode is OutputMode.AUDIO_ONLY:
            return None  # We just show the slide, no persona video

        if self.image is None:
            raise ValueError("Image is required for video generation but not available.")
        width, height = self.image.size
        video_prompt = VIDEO_PROMPT
        video_neg_prompt = VIDEO_NEG_PROMPT
        num_steps = self.get_num_steps()

        # Generate video for slide
        if output_mode == OutputMode.VIDEO_AUDIO_UNSYNCED:
            slide_video_binary = await self.gen.gen_video(
                img=self.image,
                prompt=video_prompt,
                neg_prompt=video_neg_prompt,
                width=width,
                height=height,
                video_seconds=audio_duration,  # Because of rounding, this may produce more frames than asked
                steps=num_steps,
                task_id=f"{slide_number:03d}",
                wait_request=True,
                deadline=self.get_slide_deadline(slide_number),
            )
        elif audio_duration < MAX_FT_DURATION_SECS:
            slide_video_binary = await self.gen.gen_video_audio_from_img(
                img=self.image,
                audio_base64=audio_base64,
                prompt=video_prompt,
                neg_prompt=video_neg_prompt,
                width=width,
                height=height,
                steps=num_steps,
                task_id=f"{slide_number:03d}",
                deadline=self.get_slide_deadline(slide_number),
            )
        else:
            # TODO
            # slide_video_binary = await self.gen_video_audio_from_img_chunks(
            gen_video_chunked = GenVideoChunked(
                video_id=slide_number,
                gen=self.gen,
                job_path=self.job_path,
                logger=self.logger)
            slide_video_binary = await gen_video_chunked.gen_video_chunked(
                audio_path=audio_path,
                image=self.image,
                prompt=video_prompt,
                neg_prompt=video_neg_prompt,
                width=width,
                height=height,
                num_steps=num_steps,
                upscaling=self.get_config_bool("upscaling"),
                debug=self.get_config_bool("debug_image"),
                deadline=self.get_slide_deadline(slide_number),
            )

        video_frames = await get_video_frames(slide_video_binary)
        video_file_info = get_video_file_info(slide_video_binary)
        video_info = video_file_info.get("video", {})
        video_fps = video_info.get("fps", FANTASYTALKING_FPS)

        video_audio_path = f"{self.job_path}/{slide_number:03d}_persona.mp4"
        video_audio_path = await save_video_audio(
            video_content=video_frames,
            audio_path=audio_path,
            out_video_path=video_audio_path,
            fps=video_fps)

        self._log_video_info(f"[{slide_number}] Generated slide", video_audio_path)
        self.logger.info(f"[{slide_number}] Generated slide in {time.time() - t0:.3f} seconds.")

        return video_audio_path

    def _handle_slide_exception(
        self,
        slide_number: int,
        ex: Exception
    ) -> None:
        """Handle exceptions during slide processing."""
        if isinstance(ex, ServiceError):
            self.logger.error(f"[{slide_number}] Service error: {ex}")
        elif isinstance(ex, (NoRunnableContainerError, NoActiveContainerError, ServiceNotFoundError)):
            self.logger.error(f"[{slide_number}] {ex}")
        else:
            self.logger.error(f"[{slide_number}] Error ({type(ex).__name__}): {ex}")

    def _handle_package_not_found(
        self,
        ex: Exception
    ) -> None:
        self.logger.error(
            f"Package not found for {self.job_id}: {ex}. "
            "Is file encrypted or protected by password? "
            "Make it public."
        )

    async def gen_persona(
        self,
        pptx_base64: str,
    ) -> None:
        """
        Generate a video with a character (persona) going over the slides.
        """
        async with self.job_status_handler(
            extra_handlers={
                PackageNotFoundError: self._handle_package_not_found,
            }
        ):
            if not pptx_base64:
                self.logger.error("Document is required.")
                await self.save_status(JobStatus.FAILED)
                raise ValueError("Missing 'pptx_base64' in request")
            self.logger.info(f"Generating persona for slides with {bytes_to_human(len(pptx_base64))}.")

            # Save as PDF for debugging
            self.logger.info(f"Document base64 with {bytes_to_human(len(pptx_base64))}.")
            pptx_path = f"{self.job_path}/document.pptx"
            await save_base64_as_binary(pptx_path, pptx_base64)

            # Save as text for debugging
            txt_path = f"{self.job_path}/document.txt"
            async with aiofiles.open(txt_path, "w", encoding="utf-8") as file:
                presentation = Presentation(pptx_path)
                slide_ix = 0
                for slide in presentation.slides:
                    is_hidden = slide._element.get("show") == "0"
                    if not is_hidden:
                        await file.write(f"--- Slide {slide_ix + 1} ---\n")
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    await file.write(paragraph.text + "\n")
                        slide_ix += 1

            await self.save_status(JobStatus.RUNNING)

            # Estimate number of words per slide to fit into video duration
            video_duration_seconds = self.get_config_int("video_duration_seconds", 60)
            if video_duration_seconds <= 0:
                raise ValueError("video_duration_seconds must be positive")
            num_slides = get_num_slides(pptx_path)
            num_slides = max(1, num_slides)
            seconds_per_slide = video_duration_seconds / num_slides
            num_words_per_slide = estimate_num_words_from_audio_duration(seconds_per_slide)
            self.logger.info(
                f"Estimated {num_words_per_slide} words for each of the {num_slides} slides "
                f"to fit into a {video_duration_seconds} seconds video.")

            persona_size_ratio = self.get_config_float("persona_size_ratio", 0.25)
            if not (0.05 <= persona_size_ratio <= 0.5):
                raise ValueError("persona_size_ratio must be between 5% and 50%")

            await self.save_status(JobStatus.RUNNING)

            # Save as images for generating slides video
            slide_image_paths = pptx_to_images(
                pptx_path,
                output_path=self.job_path,
                width=self.width,
                height=self.height,
                logger=self.logger,
            )
            if not slide_image_paths:
                raise ValueError(f"No images extracted from PPTX {pptx_path}. Is it empty?")

            # Persona image sketch parameters
            img_prompt = IMG_PROMPT
            img_neg_prompt = IMG_NEG_PROMPT
            persona_width = int(math.ceil(self.width * persona_size_ratio))
            persona_height = int(math.ceil(self.height * persona_size_ratio))

            await self.save_status(JobStatus.RUNNING)

            # Generate transcript for slides
            slide_video_tasks = {}
            async with aiofiles.open(f"{self.job_path}/slides_transcript.jsonl", "wb") as file:
                async for line_json in self.gen.gen_slides_transcript(
                    pptx_base64=pptx_base64,
                    max_words_per_slide=num_words_per_slide,
                    task_id=self.job_id,
                ):
                    line = (json.dumps(line_json) + "\n").encode("utf-8")
                    await file.write(line)
                    await file.flush()

                    line_type = line_json.get("type", "")
                    if line_type == "persona":
                        # gender = line_json.get("gender", "unknown")
                        img_prompt = line_json.get("description", IMG_PROMPT)

                        # Generate persona image
                        self.image_task = asyncio.create_task(
                            self.gen.gen_image(
                                img_prompt,
                                neg_prompt=img_neg_prompt,
                                width=persona_width,
                                height=persona_height,
                                steps=25,  # TODO steps
                                task_id="persona_image",
                                deadline=self.get_submission_time(),
                            ))
                    elif line_type == "slide_transcript":
                        slide_number = line_json.get("slide_number", -1)
                        slide_text = line_json.get("transcript", "")
                        MAX_LOG_TEXT_LENGTH = 60
                        self.logger.info(f"Slide {slide_number}: {slide_text[:MAX_LOG_TEXT_LENGTH]}...")

                        # Generate slide video
                        slide_video_task = asyncio.create_task(
                            self.gen_slide_video(
                                slide_number,
                                slide_text,
                            ))
                        slide_video_tasks[slide_number] = slide_video_task
                    else:
                        self.logger.info(f"Unknown line type: {line_type}")

            await self.save_status(JobStatus.RUNNING)

            # Collect slide videos
            slide_video_paths: List[Optional[str]] = []
            results = await asyncio.gather(*slide_video_tasks.values(), return_exceptions=True)
            for slide_number, result in enumerate(results):
                if isinstance(result, Exception):
                    self._handle_slide_exception(slide_number, result)
                    slide_video_paths.append(None)
                elif not result:
                    self.logger.warning(f"[{slide_number}] No video generated. Skipping...")
                    slide_video_paths.append(None)
                elif isinstance(result, str):
                    slide_video_paths.append(result)
                    self.logger.info(f"[{slide_number}] Video+audio saved to '{result}'.")
                else:
                    self.logger.warning(f"[{slide_number}] Wrong result generated {type(result)}...")
                    slide_video_paths.append(None)

            if not slide_video_paths:
                raise ValueError("No slide videos generated. Cannot create final slides video.")

            await self.save_status(JobStatus.RUNNING)

            # Concatenate all slide videos into a final slides video
            self.logger.info(
                f"Overlaying {len(slide_image_paths)} slides with "
                f"{len(slide_video_paths)} persona videos...")

            for slide_num, slide_video_path in enumerate(slide_video_paths):
                audio_path = f"{self.job_path}/{slide_num + 1:03d}.wav"
                if not await aiofiles.os.path.exists(audio_path):
                    # TODO should we continue with empty audio?
                    raise FileNotFoundError(f"Audio file for slide {slide_num} not found at '{audio_path}'.")

                if slide_num < len(slide_image_paths):
                    slide_image_path = slide_image_paths[slide_num]
                    slide_image = Image.open(slide_image_path)
                    slide_image = slide_image.resize((self.width, self.height))
                else:
                    self.logger.warning(f"Slide image for slide {slide_num} not found, generating blank slide.")
                    slide_image = get_frame_with_text(
                        width=self.width,
                        height=self.height,
                        text=f"Slide {slide_num + 1} not available.",
                        output_type="pil",
                        background_color="white",
                        font_color="black",
                    )

                if slide_video_path:
                    self.logger.info(f"[{slide_num}] Overlaying persona video on slide.")

                    slide_video_binary = await read_file_bytes(slide_video_path)
                    slide_video_frames = await get_video_frames(slide_video_binary)
                    for frame_index, slide_video_frame in enumerate(slide_video_frames):
                        slide_video_frames[frame_index] = overlay_image_on_image(
                            base_image=slide_image,
                            overlay_image=slide_video_frame,
                            position=("bottom", "right"),
                            overlay_percentage=persona_size_ratio)

                    video_file_info = get_video_file_info(slide_video_binary)
                    video_info = video_file_info.get("video", {})
                    fps = video_info.get("fps", FANTASYTALKING_FPS)
                else:
                    fps = FANTASYTALKING_FPS  # By default, use FT fps
                    audio_info = get_audio_file_info(audio_path)
                    audio_duration = audio_info["duration_seconds"]
                    num_frames = int(math.ceil(fps * audio_duration))
                    if self.image:
                        self.logger.info(
                            f"[{slide_num}] No persona video, generating static video with "
                            f"{audio_duration:.3f} seconds and {num_frames} frames.")
                        slide_image_persona = overlay_image_on_image(
                            base_image=slide_image,
                            overlay_image=self.image,
                            position=("bottom", "right"),
                            overlay_percentage=persona_size_ratio)
                        slide_video_frames = [slide_image_persona] * num_frames
                    else:
                        self.logger.info(
                            f"[{slide_num}] No persona image, generating slide video with "
                            f"{audio_duration:.3f} seconds and {num_frames} frames.")
                        slide_video_frames = [slide_image] * num_frames

                slide_video_path = f"{self.job_path}/{slide_num + 1:03d}.mp4"
                slide_video_path = await save_video_audio(
                    video_content=slide_video_frames,
                    audio_path=audio_path,
                    out_video_path=slide_video_path,
                    fps=fps)
                slide_video_paths[slide_num] = slide_video_path
                self.logger.info(f"[{slide_num}] Slide video+audio saved to '{slide_video_path}'.")

            video_binary = await concatenate_videos(slide_video_paths)
            video_path = f"{self.job_path}/{self.job_id}.mp4"
            async with aiofiles.open(video_path, "wb") as file:
                await file.write(video_binary)

    def get_slide_deadline(
        self,
        slide_number: int,
    ) -> Optional[float]:
        """
        Get the deadline for a slide.
        """
        submission_time = self.get_submission_time()
        SECONDS_PER_SLIDE = 5.0  # TODO
        slide_deadline = submission_time + (slide_number * SECONDS_PER_SLIDE)
        return slide_deadline
