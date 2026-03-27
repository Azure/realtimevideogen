#!/usr/bin/env python3
"""
Unit tests for StreamPersona.
"""

import os
import sys
import asyncio
import pytest
import aiofiles

from typing import List
from typing import Any

from PIL import Image

from http import HTTPStatus

from pptx import Presentation

from quart import Quart

from unittest.mock import patch
from unittest.mock import MagicMock
from unittest.mock import AsyncMock

# Add current path
sys.path.append(os.getcwd())

from file_utils import read_file_base64
from file_utils import read_file_bytes
from file_utils import binary_to_base64

from tests.test_utils import temp_sys_path
from tests.torch_mock import TorchMock
from tests.streamwise_app.app_test_helpers import check_app_root
from tests.streamwise_app.app_test_helpers import check_health
from tests.streamwise_app.app_test_helpers import check_files
from tests.streamwise_app.app_test_helpers import check_unknown_route
from tests.streamwise_app.app_test_helpers import check_job_submit_page
from tests.streamwise_app.app_test_helpers import check_job_status_page
from tests.streamwise_app.app_test_helpers import check_api_job_status
from tests.streamwise_app.app_test_helpers import check_api_job_requests

mock_torch = TorchMock()
mock_modules = {}
mock_modules.update(mock_torch.get_sub_modules())

with patch.dict(sys.modules, mock_modules):
    from media_utils import save_video_frames

    with temp_sys_path("apps", "apps/streampersona"):
        from apps.streampersona.streampersona import StreamPersonaApp
        from apps.streampersona.streampersona_job import StreamPersonaJob
        from apps.streampersona.streampersona_job import overlay_image_on_image
        from apps.streamwise_job import OutputMode
        from apps.client import ServiceRequest
        from apps.client import RequestStatus

streampersona_app = StreamPersonaApp()


@pytest.fixture(name="test_app")
def _test_app() -> Quart:
    return streampersona_app.app


@pytest.mark.asyncio
async def test_app(test_app: Quart) -> None:
    """Check that GET / returns 200."""
    await check_app_root(test_app, "StreamPersona")


@pytest.mark.asyncio
async def test_health(test_app: Quart) -> None:
    """Check /health."""
    await check_health(test_app)


@pytest.mark.asyncio
async def test_files(test_app: Quart) -> None:
    """Check /files endpoint."""
    await check_files(test_app, "streampersona")


@pytest.mark.asyncio
async def test_unknown_route(test_app: Quart) -> None:
    """Check that an unknown route returns 404."""
    await check_unknown_route(test_app)


@pytest.mark.asyncio
async def test_job_submit_page(test_app: Quart) -> None:
    """Check the web page for job submission."""
    await check_job_submit_page(test_app)


@pytest.mark.asyncio
async def test_job_status_page(test_app: Quart) -> None:
    """Check the web page for job status."""
    await check_job_status_page(test_app)


@pytest.mark.asyncio
async def test_api_job_status(test_app: Quart) -> None:
    """Check the API for job status (returns UNKNOWN for nonexistent jobs)."""
    await check_api_job_status(test_app)


@pytest.mark.asyncio
async def test_api_job_requests(test_app: Quart) -> None:
    """Check the API for job requests listing (returns empty for nonexistent jobs)."""
    await check_api_job_requests(test_app)


def create_example_slides(
    texts: List[str] = ["Hello", "World"],
    file_name: str = "/tmp/example.pptx"
) -> None:
    ret = Presentation()
    layout = ret.slide_layouts[0]

    for text in texts:
        slide = ret.slides.add_slide(layout)
        slide.shapes.title.text = text

    # Save the file
    ret.save(file_name)


@pytest.mark.asyncio
async def test_api_job_submit(test_app: Quart) -> None:
    """Check the API for job requests."""
    client = test_app.test_client()

    # Mock the service manager
    streampersona_app.service_manager = MagicMock()

    response = await client.post("/api/job", json={"pdf_base64": "AAAA"})
    assert response.status_code == HTTPStatus.BAD_REQUEST
    response_json = await response.get_json()
    assert "error" in response_json
    assert response_json["error"] == "Missing 'pptx_base64' in job config"

    # Success case
    create_example_slides(
        texts=["Slide 1", "Slide 2"],
        file_name="/tmp/example.pptx"
    )
    async with aiofiles.open("/tmp/example.pptx", "rb") as file:
        pptx_binary = await file.read()
        pptx_base64 = binary_to_base64(pptx_binary)

    response = await client.post("/api/job", json={"pptx_base64": pptx_base64})
    assert response.status_code == HTTPStatus.INTERNAL_SERVER_ERROR
    response_json = await response.get_json()
    assert "error" in response_json
    assert "No such file or directory: 'libreoffice'" in response_json["error"]

    """
    # TODO mimic libreoffice being present
    assert response.status_code == HTTPStatus.OK
    response_json = await response.get_json()
    assert "job_id" in response_json
    assert response_json["status"] == "success"
    """


async def _mock_generation(
    job: StreamPersonaJob
) -> StreamPersonaJob:
    """
    Mock multi-modal generation components.
    """
    audio_base64 = await read_file_base64("tests/data/audio_4675.wav")
    job.gen.service_manager.get_service_url = MagicMock(
        return_value="http://mock_service_url:1234"
    )
    job.gen.gen_audio = AsyncMock(return_value=audio_base64)
    loop = asyncio.get_running_loop()
    image_future: asyncio.Future[Image.Image] = loop.create_future()
    image_future.set_result(Image.new("RGB", (640, 480), color="blue"))
    job.image_task = image_future  # type: ignore[assignment]

    video_frames: List[Image.Image] = [
        Image.new("RGB", (640, 480), color=color)
        for color in ["blue", "green", "red", "yellow"]
    ]
    video_path = await save_video_frames(video_frames)
    video_binary = await read_file_bytes(video_path)
    job.gen.gen_video_audio_from_img = AsyncMock(return_value=video_binary)

    async def gen_video_mock(*args: Any, **kwargs: Any) -> Any:
        if kwargs.get("wait_request", None):
            return video_binary
        req = ServiceRequest(
            request_id="mock_request_id",
            service_name="hunyuanframepackf1",
            payload_json={"job_id": job.job_id},
        )
        req.status = RequestStatus.COMPLETED
        req.future = asyncio.Future()
        req.future.set_result(("video/mp4", video_binary))
        return req

    job.gen.gen_video = AsyncMock(side_effect=gen_video_mock)
    return job


@pytest.mark.asyncio
async def test_gen_slide_video() -> None:
    service_manager = MagicMock()
    job_id = "test_gen_slide_video"

    job = None
    try:
        job = StreamPersonaJob(job_id, service_manager)

        with pytest.raises(Exception, match="Error generating audio"):
            await job.gen_slide_video(
                slide_number=0,
                slide_text="Test Slide 0",
            )

        # Mocking generation
        job = await _mock_generation(job)

        await job.gen_slide_video(
            slide_number=1,
            slide_text="Test Slide 1",
        )
    finally:
        if job:
            await job.close()


@pytest.mark.asyncio
async def test_gen_slide_video_unsynced() -> None:
    service_manager = MagicMock()
    job_id = "test_gen_slide_video_unsynced"

    job = None
    try:
        job = StreamPersonaJob(job_id, service_manager)
        job.config["output_mode"] = OutputMode.VIDEO_AUDIO_UNSYNCED

        with pytest.raises(Exception, match="Error generating audio"):
            await job.gen_slide_video(
                slide_number=0,
                slide_text="Test Slide 0",
            )
        job = await _mock_generation(job)
        await job.gen_slide_video(
            slide_number=1,
            slide_text="Test Slide Unsynced",
        )
    finally:
        if job:
            await job.close()


def test_overlay_image_on_image() -> None:
    """Test overlay_image_on_image()."""
    base_image = Image.new("RGB", (200, 200), color="blue")
    overlay_image = Image.new("RGBA", (100, 100), color="pink")

    out_image = overlay_image_on_image(
        base_image,
        overlay_image,
        position=("bottom", "left")
    )
    assert out_image is not None
    assert out_image.size == base_image.size

    # Corner case: overlay larger than base
    out_image = overlay_image_on_image(
        overlay_image,
        base_image,
        position=("top", "center")
    )
    assert out_image is not None
    assert out_image.size == overlay_image.size

    # More positions
    for pos_vertical in ["top", "center", "bottom"]:
        for pos_horizontal in ["left", "center", "right"]:
            out_image = overlay_image_on_image(
                base_image,
                overlay_image,
                position=(pos_vertical, pos_horizontal)
            )
            assert out_image is not None
            assert out_image.size == base_image.size
